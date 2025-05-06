from flask import Flask, request, jsonify, send_file, send_from_directory
from flask_cors import CORS
import os
from dotenv import load_dotenv
import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tempfile
import shutil
import json
import re
import base64
from io import BytesIO
from PIL import Image
from pptx.enum.text import MSO_VERTICAL_ANCHOR

# Load environment variables
load_dotenv()

app = Flask(__name__, static_folder='sliding')
CORS(app)

# Configure OpenAI
openai.api_key = os.getenv('OPENAI_API_KEY')
if not openai.api_key:
    print("Warning: OPENAI_API_KEY not found in environment variables")

UPLOAD_FOLDER = 'uploads'

@app.route('/')
def serve_index():
    return send_from_directory(app.static_folder, 'index.html')

@app.route('/<path:path>')
def serve_static(path):
    return send_from_directory(app.static_folder, path)

@app.route('/api/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if file and file.filename.endswith(('.ppt', '.pptx')):
        filename = os.path.join(UPLOAD_FOLDER, file.filename)
        file.save(filename)
        return jsonify({'message': 'File uploaded successfully', 'filename': file.filename})
    
    return jsonify({'error': 'Invalid file type'}), 400

@app.route('/api/generate-content', methods=['POST'])
def generate_content():
    if not openai.api_key:
        return jsonify({'error': 'OpenAI API key not configured'}), 500

    data = request.json
    title = data.get('title', '')
    
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that creates concise bullet points for presentations. Each bullet point should have a short title and brief content (maximum 2 lines) separated by a colon. Never generate more than 5 bullet points."},
                {"role": "user", "content": f"Create 3-5 brief bullet points for a presentation about: {title}. Format each point as 'Title: Content'. Keep content concise and to the point."}
            ],
            max_tokens=150  # Limit response length
        )
        content = response.choices[0].message.content
        return jsonify({'content': content})
    except Exception as e:
        print(f"Error generating content: {str(e)}")
        return jsonify({'error': 'Failed to generate content'}), 500

@app.route('/api/generate-presentation', methods=['POST'])
def generate_presentation():
    try:
        data = request.json
        content = data.get('content', '')
        slide_count = data.get('slideCount', 'brief')
        
        # Generate content using OpenAI
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a presentation content generator. Create structured content for slides based on the given topic. Each slide should have a title and 3-5 bullet points. Format each point as 'Title: Content'."},
                {"role": "user", "content": f"Create content for a {slide_count} presentation about: {content}. Include {slide_count} slides with 3-5 bullet points each."}
            ],
            max_tokens=1000
        )
        
        generated_content = response.choices[0].message.content
        
        # Process the generated content into slides
        slides = []
        current_slide = {"title": "", "content": []}
        
        for line in generated_content.split('\n'):
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('Slide') or line.startswith('#'):
                if current_slide["title"]:
                    slides.append(current_slide)
                current_slide = {"title": line, "content": []}
            elif ':' in line:
                title, content = line.split(':', 1)
                current_slide["content"].append({
                    "title": title.strip(),
                    "content": content.strip()
                })
        
        if current_slide["title"]:
            slides.append(current_slide)
        
        return jsonify({"slides": slides})
        
    except Exception as e:
        print(f"Error generating presentation: {str(e)}")
        return jsonify({'error': f'Failed to generate presentation: {str(e)}'}), 500

def rgb_string_to_tuple(rgb_str, fallback):
    try:
        parts = [int(x.strip()) for x in re.split('[, ]+', rgb_str) if x.strip()]
        if len(parts) == 3:
            return tuple(parts)
    except Exception:
        pass
    return fallback

def add_logo_to_slide(slide, logo_data_url, logo_position):
    if not logo_data_url or not logo_position:
        return
    try:
        # Parse base64 data
        if ',' in logo_data_url:
            header, b64data = logo_data_url.split(',', 1)
        else:
            b64data = logo_data_url
        img_bytes = base64.b64decode(b64data)
        # Save to temp file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
            tmp.write(img_bytes)
            tmp.flush()
            tmp_path = tmp.name
        # Get image size for aspect ratio
        with Image.open(tmp_path) as img:
            width, height = img.size
            aspect = width / height
        logo_height = Inches(0.5)
        logo_width = Inches(0.5) * aspect
        slide_width = Inches(16)
        slide_height = Inches(9)
        # Position
        if logo_position == 'top-left':
            left = Inches(0.3)
            top = Inches(0.3)
        elif logo_position == 'top-right':
            left = slide_width - logo_width - Inches(0.3)
            top = Inches(0.3)
        elif logo_position == 'bottom-left':
            left = Inches(0.3)
            top = slide_height - logo_height - Inches(0.3)
        elif logo_position == 'bottom-right':
            left = slide_width - logo_width - Inches(0.3)
            top = slide_height - logo_height - Inches(0.3)
        else:
            left = Inches(0.3)
            top = Inches(0.3)
        slide.shapes.add_picture(tmp_path, left, top, logo_width, logo_height)
        os.remove(tmp_path)
    except Exception as e:
        print(f"Error adding logo: {e}")

@app.route('/api/save-presentation', methods=['POST'])
def save_presentation():
    try:
        data = request.json
        slides = data.get('slides', [])
        filename = data.get('filename', 'presentation.pptx')
        styleColors = data.get('styleColors', {})
        logoSettings = data.get('logoSettings', {})
        logo_data_url = logoSettings.get('logoDataUrl')
        logo_position = logoSettings.get('logoPosition')
        # Parse colors or use defaults
        content_text_rgb = rgb_string_to_tuple(styleColors.get('contentTextColorRGB', ''), (34,34,34))
        highlight_rgb = rgb_string_to_tuple(styleColors.get('highlightColorRGB', ''), (220,53,69))
        forms_bg_rgb = rgb_string_to_tuple(styleColors.get('formsBgColorRGB', ''), (244,246,251))
        uploaded_path = os.path.join(UPLOAD_FOLDER, filename)
        
        if os.path.exists(uploaded_path):
            prs = Presentation(uploaded_path)
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            add_title_slide = False
            # Get number of existing slides in the uploaded file
            num_existing_slides = len(prs.slides)
        else:
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            add_title_slide = True
            num_existing_slides = 0
        
        if add_title_slide:
            title_slide_layout = prs.slide_layouts[6]  # Blank layout
            title_slide = prs.slides.add_slide(title_slide_layout)
            title_box = title_slide.shapes.add_shape(
                1,  # Rectangle shape
                Inches(0.61),  # Left position
                Inches(2.16),  # Top position
                Inches(8.7),   # Width
                Inches(1.46)   # Height
            )
            title_box.fill.background()  # No fill
            title_box.line.fill.background()  # No border
            title_box.shadow.inherit = False  # Remove shadow
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.alignment = PP_ALIGN.LEFT
            title_run = title_paragraph.add_run()
            title_run.text = filename.replace('.pptx', '').title()  # Capitalize every word
            title_run.font.name = 'Frutiger 45 Light'
            title_run.font.size = Pt(40)
            title_run.font.color.rgb = RGBColor(0, 0, 0)
            title_run.font.shadow = None  # Remove shadow
            presenter_box = title_slide.shapes.add_shape(
                1,  # Rectangle shape
                Inches(0.61),  # Left position
                Inches(4.72),  # Top position
                Inches(7.48),  # Width
                Inches(0.3)    # Height
            )
            presenter_box.fill.background()  # No fill
            presenter_box.line.fill.background()  # No border
            presenter_box.shadow.inherit = False  # Remove shadow
            presenter_frame = presenter_box.text_frame
            presenter_frame.word_wrap = True
            presenter_paragraph = presenter_frame.paragraphs[0]
            presenter_paragraph.alignment = PP_ALIGN.LEFT
            presenter_run = presenter_paragraph.add_run()
            presenter_run.text = "Presenter"
            presenter_run.font.name = 'Frutiger 45 Light'
            presenter_run.font.size = Pt(16)
            presenter_run.font.color.rgb = RGBColor(0, 0, 0)
            presenter_run.font.shadow = None  # Remove shadow
        
        # Only add new slides that are not already present in the uploaded file
        for idx, slide_data in enumerate(slides[num_existing_slides:]):
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            layout_type = slide_data.get('layoutType', 'boxes')
            title_text = slide_data.get('title', '').title()
            content = slide_data.get('content', [])
            if isinstance(content, str):
                content = [line.strip() for line in content.split('\n') if line.strip()]
                content = [{'title': line.split(':')[0].strip(), 'content': line.split(':')[1].strip() if ':' in line else ''} for line in content]
                content = content[:5]

            # Add title at the top for all layouts
            title_shape = slide.shapes.add_shape(
                1,  # Rectangle shape
                Inches(1),  # Left position
                Inches(0.5),  # Top position
                Inches(14),  # Width
                Inches(1)  # Height
            )
            title_shape.fill.background()  # No fill
            title_shape.line.fill.background()  # No border
            title_shape.shadow.inherit = False  # Remove shadow
            title_frame = title_shape.text_frame
            title_frame.word_wrap = True
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.alignment = PP_ALIGN.LEFT
            title_run = title_paragraph.add_run()
            title_run.text = title_text
            title_run.font.name = 'Frutiger 45 Light'
            title_run.font.size = Pt(44)
            title_run.font.color.rgb = RGBColor(0,0,0)
            title_run.font.shadow = None  # Remove shadow

            if layout_type == 'boxes':
                create_list_boxes(slide, content, content_text_rgb, highlight_rgb, forms_bg_rgb)
            elif layout_type == 'versus':
                create_versus_layout(slide, content, content_text_rgb, highlight_rgb, forms_bg_rgb)
            elif layout_type == 'brain':
                create_brain_layout(slide, content, content_text_rgb, highlight_rgb, forms_bg_rgb)
            else:
                create_list_boxes(slide, content, content_text_rgb, highlight_rgb, forms_bg_rgb)
            # Add logo to all slides except the first
            add_logo_to_slide(slide, logo_data_url, logo_position)
        
        if not os.path.exists(UPLOAD_FOLDER):
            os.makedirs(UPLOAD_FOLDER)
        output_path = os.path.join(UPLOAD_FOLDER, 'presentation_edited.pptx')
        prs.save(output_path)
        if not os.path.exists(output_path):
            raise Exception("Failed to create presentation file")
        return send_from_directory(UPLOAD_FOLDER, 'presentation_edited.pptx', as_attachment=True)
    except Exception as e:
        print(f"Error saving presentation: {str(e)}")
        return jsonify({'error': f'Failed to save presentation: {str(e)}'}), 500

def create_list_boxes(slide, content, content_text_rgb, highlight_rgb, forms_bg_rgb):
    """Create boxes for list items."""
    # Filter out any empty content items
    content = [item for item in content if item and (isinstance(item, dict) and (item.get('title') or item.get('content')) or isinstance(item, str) and item.strip())]
    
    num_points = len(content)
    if num_points > 0:
        # Get the actual slide width from the presentation
        slide_width = Inches(16)  # Standard 16:9 slide width
        slide_margin = Inches(1)  # Margin from edges
        
        # Calculate available width for boxes
        available_width = slide_width - (2 * slide_margin)
        
        # Calculate box dimensions
        box_margin = Inches(0.3)  # Margin between boxes
        total_margins = box_margin * (num_points - 1)  # Total space needed for margins
        box_width = (available_width - total_margins) / num_points
        
        # Calculate starting position to center the boxes
        start_x = slide_margin
        start_y = Inches(3)  # Position below title
        
        for i, point in enumerate(content):
            # Calculate box position
            left = start_x + (i * (box_width + box_margin))
            top = start_y
            
            # Create box shape
            box = slide.shapes.add_shape(1, left, top, box_width, Inches(2.5))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*forms_bg_rgb)
            box.line.fill.background()
            box.shadow.inherit = False  # Remove shadow
            
            # Configure text frame
            text_frame = box.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.margin_top = Inches(0.2)
            text_frame.margin_bottom = Inches(0.2)
            
            # Clear any default paragraphs
            for paragraph in text_frame.paragraphs:
                p = paragraph._element
                p.getparent().remove(p)
            
            if isinstance(point, dict):
                title = point.get('title', '')
                content_val = point.get('content', '')
                
                # Add title
                p = text_frame.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = title.title()
                run.font.name = 'Frutiger 45 Light'
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(*highlight_rgb)
                run.font.bold = False
                run.font.shadow = None
                # Add line break (empty paragraph)
                if content_val:
                    p = text_frame.add_paragraph()
                    p.alignment = PP_ALIGN.LEFT
                    p.add_run().text = ''
                # Add content
                if content_val:
                    p = text_frame.add_paragraph()
                    run = p.add_run()
                    run.text = content_val
                    run.font.name = 'Frutiger 45 Light'
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor(*content_text_rgb)
                    run.font.shadow = None
            else:
                # Add simple text
                p = text_frame.add_paragraph()
                run = p.add_run()
                run.text = point
                run.font.name = 'Frutiger 45 Light'
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(*content_text_rgb)
                run.font.shadow = None

            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

def create_versus_layout(slide, content, content_text_rgb, highlight_rgb, forms_bg_rgb):
    # Add versus image in the center
    img_path = os.path.join('sliding', 'icons', 'versus.png')
    img_width = Inches(1.5)
    img_height = Inches(1.5)
    slide_width = Inches(16)
    center_x = (slide_width - img_width) / 2
    center_y = Inches(3.5)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, center_x, center_y, img_width, img_height)
    # Two side boxes
    left_box = slide.shapes.add_shape(1, Inches(1.2), Inches(2.5), Inches(5.5), Inches(3))
    right_box = slide.shapes.add_shape(1, Inches(9.3), Inches(2.5), Inches(5.5), Inches(3))
    for box, idx in zip([left_box, right_box], range(2)):
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(*forms_bg_rgb)
        box.line.fill.background()
        box.shadow.inherit = False
        text_frame = box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)
        for paragraph in text_frame.paragraphs:
            p = paragraph._element
            p.getparent().remove(p)
        if idx < len(content):
            point = content[idx]
            if isinstance(point, dict):
                title = point.get('title', '')
                box_p = text_frame.add_paragraph()
                box_p.alignment = PP_ALIGN.LEFT
                run = box_p.add_run()
                run.text = title.title()
                run.font.name = 'Frutiger 45 Light'
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(*highlight_rgb)
                run.font.bold = True
                run.font.shadow = None
                box_content = point.get('content', '')
                if box_content:
                    box_p = text_frame.add_paragraph()
                    box_p.alignment = PP_ALIGN.LEFT
                    run = box_p.add_run()
                    run.text = box_content
                    run.font.name = 'Frutiger 45 Light'
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor(*content_text_rgb)
                    run.font.shadow = None
            else:
                box_p = text_frame.add_paragraph()
                box_p.alignment = PP_ALIGN.LEFT
                run = box_p.add_run()
                run.text = str(point)
                run.font.name = 'Frutiger 45 Light'
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(*content_text_rgb)
                run.font.shadow = None

def create_brain_layout(slide, content, content_text_rgb, highlight_rgb, forms_bg_rgb):
    # Add brain image in the center
    img_path = os.path.join('sliding', 'icons', 'brain.png')
    img_width = Inches(1.5)
    img_height = Inches(1.5)
    slide_width = Inches(16)
    center_x = (slide_width - img_width) / 2
    center_y = Inches(3.5)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, center_x, center_y, img_width, img_height)
    # Two vertical boxes
    top_box = slide.shapes.add_shape(1, Inches(5.5), Inches(1.5), Inches(5), Inches(2.2))
    bottom_box = slide.shapes.add_shape(1, Inches(5.5), Inches(5.2), Inches(5), Inches(2.2))
    for box, idx in zip([top_box, bottom_box], range(2)):
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(*forms_bg_rgb)
        box.line.fill.background()
        box.shadow.inherit = False
        text_frame = box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)
        for paragraph in text_frame.paragraphs:
            p = paragraph._element
            p.getparent().remove(p)
        if idx < len(content):
            point = content[idx]
            if isinstance(point, dict):
                title = point.get('title', '')
                box_p = text_frame.add_paragraph()
                box_p.alignment = PP_ALIGN.LEFT
                run = box_p.add_run()
                run.text = title.title()
                run.font.name = 'Frutiger 45 Light'
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(*highlight_rgb)
                run.font.bold = True
                run.font.shadow = None
                box_content = point.get('content', '')
                if box_content:
                    box_p = text_frame.add_paragraph()
                    box_p.alignment = PP_ALIGN.LEFT
                    run = box_p.add_run()
                    run.text = box_content
                    run.font.name = 'Frutiger 45 Light'
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor(*content_text_rgb)
                    run.font.shadow = None
            else:
                box_p = text_frame.add_paragraph()
                box_p.alignment = PP_ALIGN.LEFT
                run = box_p.add_run()
                run.text = str(point)
                run.font.name = 'Frutiger 45 Light'
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(*content_text_rgb)
                run.font.shadow = None

            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

def create_flow_chart(slide, content):
    pass

def create_data_visualization(slide, content):
    pass

@app.route('/api/get-slides', methods=['POST'])
def get_slides():
    data = request.json
    filename = data.get('filename')
    if not filename:
        return jsonify({'error': 'No filename provided'}), 400
    path = os.path.join(UPLOAD_FOLDER, filename)
    if not os.path.exists(path):
        return jsonify({'error': 'File not found'}), 404

    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        title = ""
        content = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            if not text:
                continue
            if not title:
                title = text
            else:
                content.append(text)
        slides.append({
            "title": title,
            "content": "\n".join(content)
        })
    return jsonify({"slides": slides})

if __name__ == '__main__':
    if not openai.api_key:
        print("\nWarning: OpenAI API key not found!")
        print("Please add your API key to the .env file:")
        print("OPENAI_API_KEY=your_api_key_here\n")
    app.run(debug=True, port=5001) 